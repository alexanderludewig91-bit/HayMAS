"""
PowerPoint Generator Tool für HayMAS

Erstellt PowerPoint-Präsentationen aus Markdown-Struktur.
"""

import os
import re
from datetime import datetime
from typing import Dict, Any, List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Output-Verzeichnis
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "output")


def ensure_output_dir():
    """Stellt sicher, dass das Output-Verzeichnis existiert"""
    os.makedirs(OUTPUT_DIR, exist_ok=True)


def parse_markdown_slides(markdown_content: str) -> List[Dict[str, Any]]:
    """
    Parst Markdown-Inhalt in eine Liste von Folien.
    
    Erwartetes Format:
    # Folie 1: Titel
    Untertitel oder Beschreibung
    
    # Folie 2: Überschrift
    - Punkt 1
    - Punkt 2
    - Punkt 3
    
    Returns:
        Liste von Dicts mit 'title', 'subtitle', 'bullets'
    """
    slides = []
    
    # Split nach Folien-Headern
    # Matche: # Folie X: Titel oder ## Folie X: Titel
    pattern = r'^#{1,2}\s*Folie\s*\d+[:\s]*(.+?)$'
    
    # Aufteilen nach Folien
    parts = re.split(r'^#{1,2}\s*Folie\s*\d+[:\s]*', markdown_content, flags=re.MULTILINE)
    titles = re.findall(pattern, markdown_content, flags=re.MULTILINE)
    
    # Erste Teil ist vor der ersten Folie (ignorieren)
    contents = parts[1:] if len(parts) > 1 else []
    
    for i, (title, content) in enumerate(zip(titles, contents)):
        slide_data = {
            "title": title.strip(),
            "subtitle": "",
            "bullets": [],
            "notes": ""
        }
        
        lines = content.strip().split("\n")
        bullets = []
        other_text = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Bullet Points erkennen
            if line.startswith("- ") or line.startswith("* ") or line.startswith("• "):
                bullet_text = line[2:].strip()
                bullets.append(bullet_text)
            elif line.startswith("**") and line.endswith("**"):
                # Fetter Text als Untertitel
                slide_data["subtitle"] = line.strip("*").strip()
            elif not line.startswith("#"):
                other_text.append(line)
        
        slide_data["bullets"] = bullets
        
        # Wenn keine Bullets aber anderer Text, als Subtitle verwenden
        if not bullets and other_text and not slide_data["subtitle"]:
            slide_data["subtitle"] = " ".join(other_text[:2])  # Erste 2 Zeilen
        
        slides.append(slide_data)
    
    return slides


def create_ppt(
    markdown_content: str,
    output_filename: str = None,
    title_slide_title: str = None
) -> Dict[str, Any]:
    """
    Erstellt eine PowerPoint-Präsentation aus Markdown-Inhalt.
    
    Args:
        markdown_content: Markdown mit Folien-Struktur
        output_filename: Optional - Dateiname ohne .pptx
        title_slide_title: Optional - Haupttitel für Titelfolie
    
    Returns:
        Dict mit:
        - success: bool
        - filepath: Pfad zur erstellten Datei
        - slide_count: Anzahl erstellter Folien
        - error: Optional - Fehlermeldung
    """
    try:
        ensure_output_dir()
        
        # Folien parsen
        slides = parse_markdown_slides(markdown_content)
        
        if not slides:
            return {
                "success": False,
                "filepath": None,
                "slide_count": 0,
                "error": "Keine Folien im Markdown gefunden. Erwartet: '# Folie 1: Titel'"
            }
        
        # Neue Präsentation erstellen
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 Format
        prs.slide_height = Inches(7.5)
        
        # Farben definieren
        title_color = RGBColor(0x1a, 0x1a, 0x2e)  # Dunkles Blau
        text_color = RGBColor(0x33, 0x33, 0x33)   # Dunkelgrau
        accent_color = RGBColor(0x00, 0x7a, 0xcc) # Blau
        
        for i, slide_data in enumerate(slides):
            if i == 0:
                # Erste Folie als Titelfolie
                slide = _create_title_slide(prs, slide_data, title_color, accent_color)
            else:
                # Restliche Folien als Content-Folien
                slide = _create_content_slide(prs, slide_data, title_color, text_color)
        
        # Dateiname generieren
        if not output_filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"praesentation_{timestamp}"
        
        if not output_filename.endswith(".pptx"):
            output_filename = f"{output_filename}.pptx"
        
        filepath = os.path.join(OUTPUT_DIR, output_filename)
        
        # Speichern
        prs.save(filepath)
        
        return {
            "success": True,
            "filepath": filepath,
            "filename": output_filename,
            "slide_count": len(slides)
        }
        
    except Exception as e:
        return {
            "success": False,
            "filepath": None,
            "slide_count": 0,
            "error": str(e)
        }


def _create_title_slide(prs: Presentation, slide_data: Dict, title_color: RGBColor, accent_color: RGBColor):
    """Erstellt eine Titelfolie"""
    # Leeres Layout verwenden
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Hintergrund-Rechteck oben
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = title_color
    shape.line.fill.background()
    
    # Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.8),
        Inches(11.8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.LEFT
    
    # Untertitel
    if slide_data.get("subtitle"):
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(4.3),
            Inches(11.8), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = slide_data["subtitle"]
        p.font.size = Pt(24)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
    
    # Akzentlinie
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.75), Inches(4.1),
        Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()
    
    return slide


def _create_content_slide(prs: Presentation, slide_data: Dict, title_color: RGBColor, text_color: RGBColor):
    """Erstellt eine Content-Folie mit Titel und Bullet Points"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Header-Bereich
    from pptx.enum.shapes import MSO_SHAPE
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.3)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = title_color
    header.line.fill.background()
    
    # Titel (weiß auf dunklem Header)
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.35),
        Inches(11.8), Inches(0.8)
    )
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.LEFT
    
    # Bullet Points
    if slide_data.get("bullets"):
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.8),
            Inches(11.8), Inches(5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for j, bullet in enumerate(slide_data["bullets"]):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = text_color
            p.space_before = Pt(12)
            p.space_after = Pt(6)
    
    # Subtitle als zusätzlicher Text
    elif slide_data.get("subtitle"):
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(2),
            Inches(11.8), Inches(4)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["subtitle"]
        p.font.size = Pt(20)
        p.font.color.rgb = text_color
    
    return slide


# Tool-Definition für LLM Function Calling (OpenAI Format)
CREATE_PPT_TOOL = {
    "type": "function",
    "function": {
        "name": "create_ppt",
        "description": "Erstellt eine PowerPoint-Präsentation aus einer Markdown-Struktur mit Folien.",
        "parameters": {
            "type": "object",
            "properties": {
                "markdown_content": {
                    "type": "string",
                    "description": "Markdown-Inhalt mit Folien im Format: '# Folie 1: Titel' gefolgt von Bullet Points"
                },
                "output_filename": {
                    "type": "string",
                    "description": "Dateiname für die PPT (ohne .pptx). Falls nicht angegeben wird ein Timestamp verwendet."
                }
            },
            "required": ["markdown_content"]
        }
    }
}

# Anthropic Tool-Format
CREATE_PPT_TOOL_ANTHROPIC = {
    "name": "create_ppt",
    "description": "Erstellt eine PowerPoint-Präsentation aus einer Markdown-Struktur mit Folien.",
    "input_schema": {
        "type": "object",
        "properties": {
            "markdown_content": {
                "type": "string",
                "description": "Markdown-Inhalt mit Folien im Format: '# Folie 1: Titel' gefolgt von Bullet Points"
            },
            "output_filename": {
                "type": "string",
                "description": "Dateiname für die PPT (ohne .pptx). Falls nicht angegeben wird ein Timestamp verwendet."
            }
        },
        "required": ["markdown_content"]
    }
}
